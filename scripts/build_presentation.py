"""Build the dissertation presentation (<=20 slides) as an editable .pptx.

Embeds the actual experiment figures from docs/figures and uses only verified
numbers from the experiment CSVs / docs. Reproducible:

  python scripts/build_presentation.py
  -> presentation/anxiety_health_anxiety.pptx
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "figures"
OUT = ROOT / "presentation" / "anxiety_health_anxiety.pptx"

# theme
NAVY = RGBColor(0x0B, 0x2B, 0x45)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
AMBER = RGBColor(0xB8, 0x7A, 0x00)
INK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF2, 0xF6)

W, H = 13.333, 7.5


def _fig(name: str) -> str:
    p = FIG / name
    return str(p) if p.exists() else ""


def add_title_bar(slide, title, accent=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(1.02))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    bar.shadow.inherit = False
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.02), Inches(W), Inches(0.07))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent; strip.line.fill.background()
    strip.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(W - 0.9), Inches(0.74))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    r = p.runs[0]; r.font.size = Pt(25); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = "Segoe UI"


def add_bullets(slide, bullets, left, top, width, height, size=16, color=INK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        lvl, text, col = 0, b, color
        if isinstance(b, tuple):
            if len(b) == 3:
                lvl, text, col = b
            else:
                lvl, text = b
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(8)
        run = p.add_run(); run.text = ("• " if lvl == 0 else "– ") + text
        run.font.size = Pt(size - 2 * lvl); run.font.color.rgb = col; run.font.name = "Segoe UI"


def add_image_fit(slide, path, left, top, box_w, box_h, caption=None):
    if not path:
        return
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    w = box_w; h = w / aspect
    if h > box_h:
        h = box_h; w = h * aspect
    lx = left + (box_w - w) / 2
    ty = top + (box_h - h) / 2
    slide.shapes.add_picture(path, Inches(lx), Inches(ty), width=Inches(w), height=Inches(h))
    if caption:
        cb = slide.shapes.add_textbox(Inches(left), Inches(top + box_h - 0.02), Inches(box_w), Inches(0.3))
        cf = cb.text_frame; cf.word_wrap = True
        p = cf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption; r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GREY


def add_footer(slide, idx, total):
    tb = slide.shapes.add_textbox(Inches(W - 1.6), Inches(H - 0.42), Inches(1.4), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{idx} / {total}"; r.font.size = Pt(9); r.font.color.rgb = GREY


def add_table(slide, data, left, top, width, height, header_fill=NAVY, font=12):
    rows, cols = len(data), len(data[0])
    gfx = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = gfx.table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(font); run.font.name = "Segoe UI"
            if r == 0:
                run.font.bold = True; run.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                run.font.color.rgb = INK
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
    return tbl


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    slides = []

    def new(title=None, accent=ACCENT):
        s = prs.slides.add_slide(blank)
        if title:
            add_title_bar(s, title, accent)
        slides.append(s)
        return s

    CONTENT_TOP = 1.35

    # ---- 1. Title ----
    s = new()
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background(); band.shadow.inherit = False
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), Inches(W), Inches(0.06))
    strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT; strip.line.fill.background(); strip.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(W - 1.8), Inches(2.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "Detecting Anxiety and Health Anxiety in Reddit Text"
    r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Segoe UI"
    p2 = tf.add_paragraph(); r2 = p2.add_run()
    r2.text = "Weak supervision · clinical-feature architecture surgery · LLM baselines"
    r2.font.size = Pt(19); r2.font.color.rgb = RGBColor(0xAE, 0xCB, 0xE0); r2.font.name = "Segoe UI"
    tb3 = s.shapes.add_textbox(Inches(0.9), Inches(4.8), Inches(W - 1.8), Inches(1.2))
    tf3 = tb3.text_frame
    for line, sz in [("MSc Dissertation — 2026", 16),
                     ("743,879 Reddit posts · 38 subreddits · 12 experiments · 171 tests", 13)]:
        pp = tf3.paragraphs[0] if not tf3.paragraphs[0].runs else tf3.add_paragraph()
        rr = pp.add_run(); rr.text = line; rr.font.size = Pt(sz)
        rr.font.color.rgb = WHITE if sz > 14 else RGBColor(0xAE, 0xCB, 0xE0); rr.font.name = "Segoe UI"

    # ---- 2. The problem ----
    new("The problem")
    add_bullets(slides[-1], [
        "Health anxiety (illness anxiety disorder): persistent, excessive worry about having a "
        "serious illness — clinically distinct from, and far less studied than, general anxiety or depression.",
        "Reddit is a large-scale, naturalistic source of mental-health language: people describe symptoms, "
        "fears and help-seeking in their own words.",
        "Research question: can we detect anxiety — and specifically HEALTH anxiety — from text, "
        "rigorously, with clinical grounding, and honest about what actually works?",
        "Constraint that shapes everything: no clinician-labelled health-anxiety data exists, so we must "
        "build from weak / self-reported signals.",
    ], 0.6, CONTENT_TOP, W - 1.2, 5.4, size=18)

    # ---- 3. Literature & gap ----
    new("What the field has done — and the gap we target")
    add_bullets(slides[-1], [
        "Only published health-anxiety baseline: Low et al. 2020 (Reddit Mental Health Dataset) — "
        "linear SGD-L1, weighted-F1 = 0.851. Strong encoders (MentalBERT/RoBERTa, MentaLLaMA) and LLMs exist.",
        "The gaps: (1) no public CLINICIAN-labelled health-anxiety dataset; (2) subreddit-as-label is "
        "circular; (3) cross-corpus F1 collapse is the field norm (Harrigian et al.).",
        "Unclaimed territory we target:",
        (1, "Fuse clinical-instrument (SHAI) features directly INTO a transformer."),
        (1, "Hierarchical user-level modelling of a person's post stream."),
        (1, "Generative-LLM baselines (zero-shot + QLoRA) for health anxiety."),
    ], 0.6, CONTENT_TOP, W - 1.2, 5.4, size=17)

    # ---- 4. Data & protocol ----
    new("Data & evaluation protocol")
    add_bullets(slides[-1], [
        "743,879 posts across 38 subreddits: anxiety / health-anxiety / comorbid / depression / "
        "suicidality + subreddit-matched neutral controls.",
        "Weak-label positives are severely imbalanced: 69,010 anxiety · 10,659 depression · "
        "2,544 health-anxiety · 458 suicidality.",
        "Author-disjoint splits throughout — no user appears in both train and test. This is the "
        "rigorous protocol that exposes proxy-label leakage (and many papers skip it).",
    ], 0.55, CONTENT_TOP, 6.0, 5.2, size=15.5)
    add_image_fit(slides[-1], _fig("corpus_overview.png"), 6.7, CONTENT_TOP, 6.2, 5.2)

    # ---- 5. Labeling ----
    new("Labelling: weak supervision + self-disclosure")
    add_bullets(slides[-1], [
        "Tier-1 WEAK labels = subreddit-as-proxy prior ⊕ clinical lexicons (50/50). Lexicons grounded in "
        "GAD-7, SHAI / HAI, PHQ-9 and Columbia C-SSRS instruments.",
        "Self-disclosure TEST set: regex-verified “I was diagnosed with…” users (with negation/"
        "third-party filters) + subreddit-matched never-disclosed controls.",
        "eRisk protocol: the disclosure post is MASKED at test time; we score the user's remaining posts — "
        "no trivial keyword leakage.",
        "Honest framing: these are weak / proxy labels, NOT clinician annotations — every claim is bounded "
        "by that.",
    ], 0.6, CONTENT_TOP, W - 1.2, 5.4, size=16.5)

    # ---- 6. Methodology ----
    new("Methodology — model zoo & evaluation")
    add_bullets(slides[-1], [
        "Models: TF-IDF+LogReg · XGBoost (26 linguistic features) · MentalRoBERTa (single + multi-task) · "
        "DANN · Fusion (novel) · Hierarchical user model · LLMs (Qwen / MentaLLaMA zero-shot + QLoRA).",
        "Evaluation goes well beyond a single F1:",
        (1, "Author-disjoint F1 / AUROC / AUPRC with bootstrap CIs."),
        (1, "Cross-corpus zero-shot transfer (RMHD, expert-labelled ANGST)."),
        (1, "Calibration · per-subreddit thresholds · McNemar & paired-bootstrap significance."),
        (1, "SHAP interpretability · eRisk early-detection · adversarial robustness · fairness audit."),
        "12 experiments + supporting analyses; 171 automated tests.",
    ], 0.6, CONTENT_TOP, W - 1.2, 5.4, size=15.5)

    # ---- 7. Headline result ----
    new("Headline: beating the only health-anxiety baseline", GREEN)
    add_bullets(slides[-1], [
        "Task: r/HealthAnxiety vs r/Anxiety (submissions, author-disjoint) — is health-anxiety language "
        "separable from general anxiety?",
        "Weighted-F1: TF-IDF 0.886 · MentalRoBERTa 0.906 · RoBERTa-large 0.916.",
        "All beat Low et al. 2020 (0.851) — a +0.035 to +0.065 improvement over the only published "
        "baseline. This is our primary contribution.",
    ], 0.55, CONTENT_TOP, 5.6, 5.2, size=15.5)
    add_image_fit(slides[-1], _fig("stronger_models.png"), 6.3, CONTENT_TOP, 6.6, 5.2)

    # ---- 8. Markers ----
    new("What separates health anxiety from general anxiety?")
    add_bullets(slides[-1], [
        "Health-anxiety markers have large effect sizes: illness/bodily terms (Cohen's d ≈ +7.5), "
        "symptom-googling, reassurance-seeking — all SHAI-consistent.",
        "Discriminative n-grams: “my health”, “google”, “symptoms”, "
        "“results”, “disease”, “spiral” (HA) vs “work”, "
        "“job”, “medication”, “sleep” (anxiety).",
        "Negative sentiment rises monotonically with severity (anxiety → HA → depression → suicidality).",
    ], 0.55, CONTENT_TOP, 5.5, 5.2, size=15)
    add_image_fit(slides[-1], _fig("exp4__marker_heatmap.png"), 6.2, CONTENT_TOP, 6.7, 5.2)

    # ---- 9. External validity ----
    new("External validity — and a transfer puzzle")
    add_bullets(slides[-1], [
        "Zero-shot transfer: RMHD (Low 2020) AUROC 0.92; ANGST (3 expert psychologists) AUROC 0.82 — "
        "genuine external validity, rare in this field.",
        "Puzzle: cheap TF-IDF OUT-transferred the transformer (RMHD 0.920 vs 0.897) — this motivates the "
        "architecture surgery on the next slide.",
        "The classic F1-collapse under shift (in-dist 0.93 → cross 0.31) is a thresholding artifact: "
        "AUROC stays ~0.99 → fixed by per-subreddit calibration.",
    ], 0.55, CONTENT_TOP, 5.7, 5.2, size=15)
    add_image_fit(slides[-1], _fig("external_validation.png"), 6.4, CONTENT_TOP, 6.5, 5.2)

    # ---- 10. Contribution 1: fusion ----
    new("Contribution 1: clinical-feature architecture surgery", GREEN)
    add_bullets(slides[-1], [
        "FusionMultiTaskModel: encoder embedding (768-d) ⊕ 26 linguistic + 7 SHAI features → fusion MLP; "
        "+ focal loss; + attention pooling (all independently ablatable).",
        "WIN — fusion+focal lifts the imbalance-limited rare classes: health-anx F1 0.508→0.559, "
        "suicidality 0.444→0.522.",
        "WIN — and cross-corpus transfer: RMHD AUROC 0.894→0.931 (now BEATS TF-IDF 0.920); ANGST →0.811.",
        "NOVEL: fusing SHAI clinical-instrument features into a transformer closes the documented transfer "
        "gap — not previously done in this domain.",
    ], 0.55, CONTENT_TOP, 5.7, 5.3, size=14.5)
    add_image_fit(slides[-1], _fig("fusion_ablation.png"), 6.35, CONTENT_TOP, 6.6, 5.2)

    # ---- 11. Fusion calibration ----
    new("Contribution 1 (cont.): calibrating the new model", GREEN)
    add_bullets(slides[-1], [
        "Focal loss makes the model UNDER-confident — every learned temperature < 1 (0.47–0.62).",
        "Temperature scaling restores excellent calibration: anxiety ECE 0.020→0.006; rarer classes "
        "≤0.002 (a 69–95% reduction).",
        "Per-subreddit thresholds lift anxiety macro-F1 0.852→0.888.",
        "Both calibration tools transfer cleanly to the fused architecture — the trade focal makes "
        "(calibration for recall/ranking) is fully recoverable post-hoc.",
    ], 0.55, CONTENT_TOP, 5.7, 5.3, size=15)
    add_image_fit(slides[-1], _fig("fusion_calibration.png"), 6.35, CONTENT_TOP, 6.6, 5.2)

    # ---- 12. Contribution 2: hierarchical (negative) ----
    new("Contribution 2: hierarchical user model (honest negative)", AMBER)
    add_bullets(slides[-1], [
        "Frozen MentalRoBERTa post-encoder → learned attention aggregation over a user's post stream → "
        "user head; evaluated user-level on the disclosure test set.",
        "Result: learned attention ≤ naive mean-pooling on all three targets; both only TIE cheap TF-IDF "
        "(~0.74 user-AUROC).",
        "The bottleneck is the noisy proxy label + subreddit-matched hard negatives — NOT the aggregation "
        "mechanism. A clean null (and a useful one).",
    ], 0.55, CONTENT_TOP, 5.7, 5.2, size=15)
    add_image_fit(slides[-1], _fig("hier_user.png"), 6.35, CONTENT_TOP, 6.6, 5.2)

    # ---- 13. Contribution 3: LLMs ----
    new("Contribution 3: generative-LLM baselines", AMBER)
    add_bullets(slides[-1], [
        "Qwen2.5-7B ZERO-SHOT: weighted-F1 0.782 — loses to TF-IDF (0.886) and to the encoders.",
        "Qwen2.5-7B QLoRA (1 epoch, 4-bit): 0.917 — statistically TIED with RoBERTa-large (0.916), at "
        "20–55× the parameters.",
        "Lesson: fine-tuning — not prompting — closes the gap; the small fine-tuned encoder stays the "
        "efficient choice.",
        "MentaLLaMA verbalizer-unsuitable (chance, both prompt formats); Llama-3.1-8B gated → deferred.",
    ], 0.55, CONTENT_TOP, 5.7, 5.3, size=14.5)
    add_image_fit(slides[-1], _fig("llm_baselines.png"), 6.35, CONTENT_TOP, 6.6, 5.2)

    # ---- 14. Contribution 4: DAPT (null) ----
    new("Contribution 4: domain-adaptive pretraining (null)", AMBER)
    add_bullets(slides[-1], [
        "Continue masked-LM pretraining of roberta-base on 200k in-domain posts, then fine-tune on the HA task.",
        "No benefit: vanilla roberta-base already MATCHES MentalRoBERTa (0.906 ≈ 0.905); light DAPT slightly "
        "hurts (0.903).",
        "Consistent lesson (DAPT + LLM): at the ~0.92 Reddit-binary ceiling, fine-tuning dominates — "
        "pretraining provenance and model scale barely matter.",
    ], 0.55, CONTENT_TOP, 5.7, 5.2, size=15)
    add_image_fit(slides[-1], _fig("dapt_mlm.png"), 6.35, CONTENT_TOP, 6.6, 5.2)

    # ---- 15. Supporting rigor ----
    new("Supporting rigour (the full evaluation menu)")
    add_bullets(slides[-1], [
        "Calibration — TF-IDF was badly under-confident (ECE 0.200→0.035, −82%); transformers already "
        "well-calibrated.",
        "Per-subreddit thresholds — anxiety macro-F1 0.719→0.781.",
        "Significance — McNemar + paired bootstrap: model differences are NOT significant (honestly, tied).",
        "eRisk early detection — anxiety flagged at a median of 1 post (recall 0.78).",
        "Plus: SHAP interpretability · robustness (<5% flips under perturbation) · fairness (TPR gaps "
        "≤0.06) · weak-label filtering (null).",
    ], 0.55, CONTENT_TOP, 5.9, 5.3, size=14.5)
    add_image_fit(slides[-1], _fig("erisk.png"), 6.5, CONTENT_TOP, 6.4, 5.2)

    # ---- 16. What is new / better ----
    new("What is new / better than the literature", GREEN)
    add_bullets(slides[-1], [
        "✔ Beat the only published health-anxiety baseline (Low 2020 0.851 → 0.916 weighted-F1).",
        "✔ NOVEL: SHAI clinical-feature fusion into a transformer → closes the cross-corpus transfer gap "
        "(no prior work does this for health anxiety).",
        "✔ First DANN application to Reddit mental-health + genuine external validation on RMHD and "
        "expert-labelled ANGST.",
        "✔ Full calibration / significance / fairness / robustness / eRisk suite — beyond typical "
        "single-metric papers.",
        "✔ Scientific honesty: documented, reproducible NEGATIVE results (hierarchical, DANN, DAPT) — "
        "rare and valuable.",
    ], 0.6, CONTENT_TOP, W - 1.2, 5.4, size=16, color=INK)

    # ---- 17. What is NOT better ----
    new("What is NOT better — honest limitations", AMBER)
    add_bullets(slides[-1], [
        "✘ At/near the ~0.92 Reddit-binary ceiling → limited headroom; absolute gains are modest.",
        "✘ Weak / proxy labels (subreddit-as-label); no clinician-annotated health-anxiety gold set.",
        "✘ User-level detection only TIES cheap TF-IDF (~0.74) — deep models give no advantage there.",
        "✘ LLMs do NOT beat fine-tuned encoders: zero-shot loses, QLoRA only reaches parity (at huge cost).",
        "✘ DAPT adds nothing here; most headline runs are single-seed (no multi-seed CI averaging yet).",
    ], 0.6, CONTENT_TOP, W - 1.2, 5.4, size=16, color=INK)

    # ---- 18. Summary table ----
    new("Summary of contributions")
    table = [
        ["Contribution", "New?", "vs literature"],
        ["HA-vs-Anxiety classifier", "extends Low 2020", "✔ +0.065 F1"],
        ["SHAI clinical-feature fusion model", "NOVEL", "✔ closes transfer gap"],
        ["Calibrating the focal model", "novel application", "✔ ECE → ≤0.006"],
        ["Hierarchical user model", "novel attempt", "✘ null (ties TF-IDF)"],
        ["DANN on Reddit mental-health", "first application", "✘ null"],
        ["LLM zero-shot / QLoRA baselines", "new for HA", "≈ parity, not better"],
        ["Domain-adaptive MLM (DAPT)", "completeness", "✘ null"],
        ["Calibration/significance/fairness/robustness/eRisk", "thorough", "✔ exceeds norm"],
    ]
    add_table(slides[-1], table, 1.0, CONTENT_TOP + 0.05, 11.3, 5.0, font=14)

    # ---- 19. Future work ----
    new("Future work")
    add_bullets(slides[-1], [
        "Llama-3.1-8B baselines (zero-shot + QLoRA) — pipeline ready; pending HF gated access.",
        "MentaLLaMA via generate-and-parse decoding (the verbalizer probe is unsuitable for a long-form "
        "domain LLM).",
        "Ordinal severity head (CORAL / CORN) on the fusion model — grade severity, not just presence.",
        "The field's biggest gap: clinician-annotated health-anxiety labels + multi-seed confidence "
        "intervals.",
        "Live user-history enrichment for user-level modelling (currently blocked by Reddit IP limits).",
    ], 0.6, CONTENT_TOP, W - 1.2, 5.4, size=17)

    # ---- 20. Conclusion ----
    new("Conclusion", GREEN)
    add_bullets(slides[-1], [
        "The winning recipe: fine-tune a compact encoder + fuse clinical (SHAI) features + focal loss → "
        "beats the published baseline AND closes the cross-corpus transfer gap.",
        "LLMs match but do not beat small fine-tuned encoders; pretraining provenance and scale barely "
        "matter at the ceiling — fine-tuning on the task does.",
        "Equal weight to honest negatives (hierarchical, DANN, DAPT) and full evaluation rigour → a "
        "credible, reproducible health-anxiety NLP study.",
        "Take-away: clinically-grounded, parameter-efficient, and honest beats big-and-vague.",
    ], 0.6, CONTENT_TOP, W - 1.2, 5.4, size=17)

    total = len(slides)
    for i, s in enumerate(slides[1:], start=2):
        add_footer(s, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}  ({total} slides)")


if __name__ == "__main__":
    main()
